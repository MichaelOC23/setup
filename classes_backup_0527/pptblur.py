from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image, ImageFilter
import io
import os

#! To install needed libraries:  pip install python-pptx Pillow



class powerpoint_tools():
    def __init__(self):
        pass
    
    def ungroup_all_shapes(self, path_to_presentation):
        
        def ungroup_shapes(slide):
            """ Ungroups all grouped shapes in a slide """
            shapes_to_process = list(slide.shapes)
            ungrouped = False

            while True:
                for shape in shapes_to_process:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        group_shape = shape
                        shapes_to_add = [shp for shp in group_shape.shapes]
                        
                        # Ungroup the group itself
                        slide.shapes._spTree.remove(group_shape._element)
                        
                        # Add ungrouped shapes back to the slide
                        for shp in shapes_to_add:
                            slide.shapes._spTree.append(shp._element)

                        ungrouped = True
                
                if not ungrouped:
                    break
                
                shapes_to_process = list(slide.shapes)
                ungrouped = False
        
        """ Main function to ungroup all grouped objects in a PowerPoint presentation """
        # Load the presentation
        prs = Presentation(path_to_presentation)
        total_slides = len(prs.slides)
        progress = 0

        for slide in prs.slides:
            ungroup_shapes(slide)
            progress += 1
            print(f'Processed slide {progress}/{total_slides}')

        # Save the modified presentation
        new_pptx_path = os.path.splitext(path_to_presentation)[0] + "_ungrouped.pptx"
        prs.save(new_pptx_path)
        print(f'All slides processed. Saved new presentation as {new_pptx_path}')

    

    def blur_all_images_in_pptx(self, file_path):
        def apply_blur_to_image(image_bytes, radius):
            image = Image.open(io.BytesIO(image_bytes))
            blurred_image = image.filter(ImageFilter.GaussianBlur(radius))
            blurred_image_bytes = io.BytesIO()
            blurred_image.save(blurred_image_bytes, format='PNG')
            return blurred_image_bytes.getvalue()
        
        
        
        presentation = Presentation(file_path)
        total_slides = len(presentation.slides)
        total_shapes = sum(len(slide.shapes) for slide in presentation.slides)
        shape_count = 0
        processed_shape_ids = set()
        for slide in presentation.slides:
            shapes_to_remove = []

            shapes_in_slide = len(slide.shapes)
            print(f'Processing slide {slide.slide_id} named: {slide.name} with {shapes_in_slide} shapes. {shape_count} of {total_shapes} shapes processed. {shape_count / total_shapes * 100:.2f}% complete.')
            
            unique_slide_id = slide.slide_id
            output_path = f'/Users/michasmi/Downloads/blur/{unique_slide_id}_blurred.pptx'
            presentation.save(output_path)
            print(f'Processed file saved as {output_path}')
            
            for shape in slide.shapes:
                shape_count += 1
                print 
                print("Processing shape with shape name: ", shape.name, " and shape type: ", shape.shape_type, f"Percent complete: {shape_count / total_shapes * 100:.2f}%", end="\r", flush=True)
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.shape_id not in processed_shape_ids:
                    print("Processing image...")
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        blurred_image_bytes = apply_blur_to_image(image_bytes, 45)
                        
                        unique_shape_id = shape._element.get('id')
                        
                        # Save the blurred image to a temporary file
                        temp_image_path = f'/Users/michasmi/Downloads/blur/{unique_shape_id}_temp_image.png'
                        with open(temp_image_path, 'wb') as f:
                            f.write(blurred_image_bytes)
                        
                        # Add the blurred image and mark the original for removal
                        new_shape = slide.shapes.add_picture(temp_image_path, shape.left, shape.top, shape.width, shape.height)
                        processed_shape_ids.add(new_shape.shape_id)
                        shapes_to_remove.append(shape)
                        
                        # Remove the temporary file
                        os.remove(temp_image_path)
                    except Exception as e:
                        print(f'Error processing shape: {e}')
                        continue
        
            # Save the modified presentation
            output_path = 'blurred_' + os.path.basename(file_path)
            presentation.save(output_path)
            print(f'Processed file saved as {output_path}')


        
ppt = powerpoint_tools()
ungrouped_path = ppt.ungroup_all_shapes("/Users/michasmi/Downloads/blur/cb1.pptx")
ppt.blur_all_images_in_pptx(ungrouped_path)
